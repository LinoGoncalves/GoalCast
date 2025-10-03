from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    (
        "GoalCast: Tooling & Workflow Review",
        [
            "Purpose: Quick review for Python Standards Working Group",
            "30-minute read + homework for next week's vetting session",
            "Presenter: Team Lead",
        ],
    ),
    (
        "TL;DR Recommendation",
        [
            "Use GoalCast scaffold as a teaching/template for best practices",
            "Adopt Feature-Branch + PR workflow, pre-commit hooks, GitHub Actions CI",
            "Present it with a Teaching Mode (simplified intro + exercises)",
        ],
    ),
    (
        "Strengths - Why this is valuable",
        [
            "Modern tooling: uv, ruff, mypy, pre-commit, pytest, bandit, semgrep",
            "Production-grade workflows: CI, pre-commit, branch protection",
            "Comprehensive docs: workflow, branching, quick-start, PR template",
        ],
    ),
    (
        "Challenges & Considerations",
        [
            "Tooling overload for juniors - recommend staged adoption",
            "Domain complexity (advanced stats/ML) - not required for tooling review",
            "uv & ruff are newer - training needed for teams used to pip/black",
        ],
    ),
    (
        "Teaching & Adoption Plan",
        [
            "Short workshop (2-3h) or Tooling Tour (1h)",
            "Create a Junior Dev 'cheat sheet' and 1–2 hands-on exercises",
            "Test install on a clean machine and prepare failure demos",
        ],
    ),
    (
        "Workshop Options",
        [
            "Option A (2-3h): Full showcase + hands-on",
            "Option B (1h): Tooling tour + live demo",
            "Option C (4h): Build your own scaffold workshop",
        ],
    ),
    (
        "Immediate Next Steps (30-min ask)",
        [
            "Ask the team to READ docs/WORKFLOW.md and docs/BRANCHING_STRATEGY.md",
            "Pick preferred workshop option and volunteers for a Standards WG",
            "Bring questions and concerns for formal vetting next week",
        ],
    ),
    (
        "Appendix - Resources",
        [
            "Links to docs in the repo: docs/ (WORKFLOW.md, BRANCHING_STRATEGY.md)",
            "Commands cheat sheet: uv run pre-commit install; uv sync; uv run pytest",
            "Contact: Team Lead / Tech Lead",
        ],
    ),
]

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

for title, bullets in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = title
    tf = body_shape.text_frame
    tf.text = bullets[0]
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

prs.save("docs/GoalCast_Tooling_Review.pptx")
print("Generated docs/GoalCast_Tooling_Review.pptx")
